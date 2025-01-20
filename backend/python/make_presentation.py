import os
import json
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
from wordcloud import WordCloud
from textblob import TextBlob
import mysql.connector  # Import MySQL connector

# Set user ID and output directories
user_id = 1  # Replace with the actual user ID obtained from session
output_dir_json = os.path.join('..\\..\\public\\static\\presentations')
presentation_name = f"user{user_id}_presentation.pptx"
presentation_path = os.path.join(output_dir_json, presentation_name)

# Function to save or update image in the database
def save_or_update_image_in_db(user_id, image_path, image_type):
    # Add database interaction code here to save or update image details
    pass

servername = "sql12.freesqldatabase.com"
username = "sql12756836"
password = "qEaH9rPgZn"
database = "sql12756836"

def get_db_connection():
    return mysql.connector.connect(
        host=servername,
        user=username,
        password=password,
        database=database
    )


def delete_existing_presentation(user_id, presentation_name):
    conn = get_db_connection()
    cursor = conn.cursor()

    # Retrieve the presentation path from the database
    cursor.execute("SELECT presentation_path FROM user_presentations WHERE user_id = %s", (user_id,))
    result = cursor.fetchone()

    if result:
        # Database stores relative path like 'static/presentations/user1_presentation.pptx'
        relative_presentation_path = result[0]
        
        # Convert relative path to absolute path for deletion
        absolute_presentation_path = os.path.join('..\\..\\public\\static\\presentations', relative_presentation_path)

        # Check if the path matches the expected presentation name
        if relative_presentation_path.endswith(presentation_name):
            try:
                os.remove(absolute_presentation_path)
                print(f"Existing presentation '{absolute_presentation_path}' has been deleted.")
                
                # Delete the entry from the database
                cursor.execute("DELETE FROM user_presentations WHERE user_id = %s", (user_id,))
                conn.commit()
                print("Entry deleted from the database.")
            except FileNotFoundError:
                print(f"File '{absolute_presentation_path}' not found in the file system.")
        else:
            print(f"Presentation path mismatch for user ID {user_id}. No deletion performed.")
    else:
        print(f"No presentation found for user ID {user_id}.")

    cursor.close()
    conn.close()

# Function to update PowerPoint presentation with image and analysis text
def update_presentation(user_id, image_path, analysis_text=None):
    # Retrieve the presentation path from the database
    conn = get_db_connection()
    cursor = conn.cursor()
    
    cursor.execute("SELECT presentation_path FROM user_presentations WHERE user_id = %s", (user_id,))
    result = cursor.fetchone()
    
    if result:
        ppt_path = result[0]
    else:
        # Create a new presentation if not found
        ppt_path = f'../../public/static/presentations/user{user_id}_presentation.pptx'
        prs = Presentation()
        prs.save(ppt_path)
        cursor.execute("INSERT INTO user_presentations (user_id, presentation_path) VALUES (%s, %s)", (user_id, ppt_path))
        conn.commit()
    
    # Open the presentation
    prs = Presentation(ppt_path)
    
    # Add image slide
    image_slide_layout = prs.slide_layouts[5]
    image_slide = prs.slides.add_slide(image_slide_layout)
    img = image_slide.shapes.add_picture(image_path, Inches(1), Inches(1), height=Inches(5))
    
    # Add analysis text if available
    if analysis_text:
        text_slide_layout = prs.slide_layouts[5]
        text_slide = prs.slides.add_slide(text_slide_layout)
        textbox = text_slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5.5))
        text_frame = textbox.text_frame
        p = text_frame.add_paragraph()
        p.text = analysis_text
        p.font.size = Inches(0.2)
    
    # Save the presentation
    prs.save(ppt_path)
    
    print(f"Presentation updated with image '{image_path}' and analysis at {ppt_path}")
    
    cursor.close()
    conn.close()

# Collect and consolidate all images and their respective analysis results
# Collect and consolidate all images and their respective analysis results
def consolidate_analysis_and_images(user_id):
    output_dir = '../../public/static/images/'
    results_file = '../../public/static/analysis_results.json'
    semantic_analysis_file = '../../config/review/semantic_analysis.json'
    consolidated_results = {}
    
    # List of images to be included in the presentation
    image_paths = [
        '../../public/static/images/mrr_by_bu.png',
        '../../public/static/images/revenue_by_product_bar_chart.png',
        '../../public/static/images/churn_rate_stacked_bar_chart.png',
        '../../public/static/images/revenue_by_product_pie_chart.png',
        # Skipped the following images
        '../../public/static/images/sentiment_distribution_bar_chart.png',
        # 'static/images/wordcloud.png',
        # 'review/sentiment_polarity_graph.png',
        '../../public/static/images/revenue_by_item_currency.png',
        '../../public/static/images/revenue_quantity_by_bu.png',
        '../../public/static/images/top_customers_by_revenue.png'
    ]
    
    # Read analysis results from JSON files
    with open(results_file, 'r') as file:
        results = json.load(file)
    
    with open(semantic_analysis_file, 'r') as file:
        semantic_analysis = json.load(file)
    
    # Consolidate all results
    consolidated_results.update(results)
    consolidated_results.update(semantic_analysis)
    
    # Save the consolidated results to a new JSON file
    consolidated_results_file = '../../public/static/consolidated_analysis_results.json'
    with open(consolidated_results_file, 'w') as file:
        json.dump(consolidated_results, file, indent=4)
    
    # Update the presentation with all images and their respective analysis
    for image_path in image_paths:
        image_name = os.path.basename(image_path).split('.')[0]
        
        # Skip analysis for the specified images
        if image_path in ['../../public/static/images/sentiment_distribution_bar_chart.png',
                          '../../public/static/images/wordcloud.png',
                          '../../config/review/sentiment_polarity_graph.png']:
            update_presentation(user_id, image_path)  # Update without analysis
        else:
            analysis_text = consolidated_results.get(image_name, "No analysis available.")
            update_presentation(user_id, image_path, analysis_text)



# Main execution
def main():
    user_id = 1  # Replace with actual user ID
    delete_existing_presentation(user_id, presentation_name)
    consolidate_analysis_and_images(user_id)
    print("Presentation generation complete.")

if __name__ == "__main__":
    main()
